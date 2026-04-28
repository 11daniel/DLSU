from pptx import Presentation

def align_slides_to_first(pptx_path, output_path):
    prs = Presentation(pptx_path)
    
    if len(prs.slides) < 2:
        print("Not enough slides to align.")
        return
    
    reference_slide = prs.slides[0]
    
    # Extract positions of elements from reference slide
    reference_shapes = []
    for shape in reference_slide.shapes:
        if hasattr(shape, "left") and hasattr(shape, "top") and hasattr(shape, "width") and hasattr(shape, "height"):
            reference_shapes.append((shape.left, shape.top, shape.width, shape.height))
    
    for slide in prs.slides[1:]:
        shapes = [shape for shape in slide.shapes if hasattr(shape, "left")]
        
        # Ensure that we have the same number of elements before applying positions
        for shape, ref_shape in zip(shapes, reference_shapes):
            shape.left, shape.top, shape.width, shape.height = ref_shape
    
    prs.save(output_path)
    print(f"Presentation formatted and saved as {output_path}")


# Example usage
input_pptx = r"C:\Users\nicos\OneDrive\Documents\DLSU\Term 8\NSSECU3\FileAnalysis\screenshots_presentation.pptx"  # Change to your directory path
output_pptx = "screenshots_presentation.pptx"
align_slides_to_first(input_pptx, output_pptx)
